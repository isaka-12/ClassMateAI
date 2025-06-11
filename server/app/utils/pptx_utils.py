from pptx import Presentation
import io

def extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text